#!/usr/bin/env python3
"""
Enhanced ChatGPT Canvas to PowerPoint Converter
Addresses all potential failure points with robust error handling and enhanced features.
"""

import os
import re
import sys
import time
import json
import logging
import tempfile
import unicodedata
from pathlib import Path
from urllib.parse import urlparse
from typing import Optional, Tuple, List, Dict, Any
from dataclasses import dataclass
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from docx import Document

# Core dependencies with error handling
try:
    from selenium import webdriver
    from selenium.webdriver.chrome.options import Options
    from selenium.webdriver.chrome.service import Service
    from selenium.webdriver.common.by import By
    from selenium.webdriver.support.ui import WebDriverWait
    from selenium.webdriver.support import expected_conditions as EC
    from selenium.common.exceptions import WebDriverException, TimeoutException
    from webdriver_manager.chrome import ChromeDriverManager
except ImportError as e:
    print(f"❌ Selenium dependencies missing: {e}")
    print("Install with: pip install selenium webdriver-manager")
    sys.exit(1)

try:
    from bs4 import BeautifulSoup, NavigableString
    from bs4.element import Tag
except ImportError as e:
    print(f"❌ BeautifulSoup missing: {e}")
    print("Install with: pip install beautifulsoup4 lxml")
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_THEME_COLOR
except ImportError as e:
    print(f"❌ python-pptx missing: {e}")
    print("Install with: pip install python-pptx")
    sys.exit(1)

try:
    import requests
    from PIL import Image
    import io
except ImportError as e:
    print(f"⚠️  Optional dependencies missing (images won't be processed): {e}")
    print("Install with: pip install requests Pillow")


# === CONFIGURATION ===
@dataclass
class Config:
    """Configuration settings for the converter"""
    max_wait_time: int = 30
    page_load_timeout: int = 45
    retry_attempts: int = 3
    max_slide_content_length: int = 1000
    max_filename_length: int = 100
    supported_image_formats: List[str] = None
    font_fallbacks: Dict[str, str] = None
    
    def __post_init__(self):
        if self.supported_image_formats is None:
            self.supported_image_formats = ['.png', '.jpg', '.jpeg', '.gif', '.bmp']
        if self.font_fallbacks is None:
            self.font_fallbacks = {
                'default': 'Calibri',
                'code': 'Courier New',
                'math': 'Cambria Math',
                'fallback': 'Arial'
            }


# === LOGGING SETUP ===
def setup_logging(log_level: str = 'INFO') -> logging.Logger:
    """Setup comprehensive logging"""
    logger = logging.getLogger('canvas_converter')
    logger.setLevel(getattr(logging, log_level.upper()))
    
    if not logger.handlers:
        handler = logging.StreamHandler()
        formatter = logging.Formatter(
            '%(asctime)s - %(name)s - %(levelname)s - %(message)s'
        )
        handler.setFormatter(formatter)
        logger.addHandler(handler)
    
    return logger


# === UTILITIES ===
class SafeFilename:
    """Utility class for safe filename handling"""
    
    @staticmethod
    def sanitize(filename: str, max_length: int = 100) -> str:
        """Create safe filename with proper sanitization"""
        # Remove or replace dangerous characters
        filename = re.sub(r'[<>:"/\\|?*]', '_', filename)
        filename = re.sub(r'\s+', '_', filename)
        
        # Handle Unicode characters
        filename = unicodedata.normalize('NFKD', filename)
        filename = ''.join(c for c in filename if ord(c) < 128)
        
        # Truncate and ensure extension
        if len(filename) > max_length - 5:  # Leave room for .pptx
            filename = filename[:max_length - 5]
        
        return filename if filename else 'canvas_export'
    
    @staticmethod
    def ensure_unique(filepath: Path) -> Path:
        """Ensure filename is unique by adding counter if needed"""
        if not filepath.exists():
            return filepath
        
        counter = 1
        stem = filepath.stem
        suffix = filepath.suffix
        parent = filepath.parent
        
        while True:
            new_path = parent / f"{stem}_{counter}{suffix}"
            if not new_path.exists():
                return new_path
            counter += 1


class WebDriverManager:
    """Enhanced WebDriver management with better error handling"""
    
    def __init__(self, config: Config, logger: logging.Logger):
        self.config = config
        self.logger = logger
        self.driver = None
    
    def create_driver(self) -> webdriver.Chrome:
        """Create Chrome WebDriver with robust configuration"""
        options = Options()
        
        # Enhanced options for reliability
        options.add_argument("--headless=new")
        options.add_argument("--no-sandbox")
        options.add_argument("--disable-dev-shm-usage")
        options.add_argument("--disable-gpu")
        options.add_argument("--disable-web-security")
        options.add_argument("--disable-features=VizDisplayCompositor")
        options.add_argument("--window-size=1920,1080")
        options.add_argument("--user-agent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36")
        
        # Performance optimizations
        options.add_argument("--disable-extensions")
        options.add_argument("--disable-plugins")
        options.add_argument("--disable-images")  # Faster loading
        
        try:
            # Use webdriver-manager for automatic ChromeDriver management
            service = Service(ChromeDriverManager().install())
            driver = webdriver.Chrome(service=service, options=options)
            driver.set_page_load_timeout(self.config.page_load_timeout)
            return driver
        except Exception as e:
            self.logger.error(f"Failed to create WebDriver: {e}")
            raise WebDriverException(f"WebDriver initialization failed: {e}")
    
    def __enter__(self):
        self.driver = self.create_driver()
        return self.driver
    
    def __exit__(self, exc_type, exc_val, exc_tb):
        if self.driver:
            try:
                self.driver.quit()
            except Exception as e:
                self.logger.warning(f"Error closing WebDriver: {e}")


class ContentExtractor:
    """Advanced content extraction with multiple strategies"""
    
    def __init__(self, config: Config, logger: logging.Logger):
        self.config = config
        self.logger = logger
    
    def validate_url(self, url: str) -> bool:
        """Validate if URL is accessible and likely a ChatGPT canvas"""
        try:
            parsed = urlparse(url)
            if not parsed.scheme or not parsed.netloc:
                return False
            
            # Basic check for ChatGPT domains
            valid_domains = ['chatgpt.com', 'chat.openai.com']
            if not any(domain in parsed.netloc for domain in valid_domains):
                self.logger.warning(f"URL doesn't appear to be a ChatGPT canvas: {url}")
            
            return True
        except Exception as e:
            self.logger.error(f"URL validation failed: {e}")
            return False
    
    def scrape_with_retry(self, url: str) -> Tuple[Optional[BeautifulSoup], Optional[Tag]]:
        """Scrape with multiple retry attempts and strategies"""
        
        if not self.validate_url(url):
            return None, None
        
        for attempt in range(self.config.retry_attempts):
            try:
                self.logger.info(f"Scraping attempt {attempt + 1}/{self.config.retry_attempts}")
                return self._scrape_single_attempt(url)
            except Exception as e:
                self.logger.warning(f"Attempt {attempt + 1} failed: {e}")
                if attempt < self.config.retry_attempts - 1:
                    time.sleep(2 ** attempt)  # Exponential backoff
                else:
                    self.logger.error("All scraping attempts failed")
        
        return None, None
    
    def _scrape_single_attempt(self, url: str) -> Tuple[Optional[BeautifulSoup], Optional[Tag]]:
        """Single scraping attempt with comprehensive error handling"""
        
        with WebDriverManager(self.config, self.logger) as driver:
            try:
                self.logger.info(f"Loading URL: {url}")
                driver.get(url)
                
                # Wait for basic page structure
                WebDriverWait(driver, self.config.max_wait_time).until(
                    EC.presence_of_element_located((By.TAG_NAME, "body"))
                )
                
                # Additional waits for dynamic content
                self._wait_for_content_load(driver)
                
                html = driver.page_source
                soup = BeautifulSoup(html, "lxml")
                
                # Multiple content extraction strategies
                content = self._extract_content_multiple_strategies(soup)
                
                if content:
                    self.logger.info("Content successfully extracted")
                    return soup, content
                else:
                    self.logger.warning("No content found with any extraction strategy")
                    return soup, None
                    
            except TimeoutException:
                self.logger.error(f"Page load timeout after {self.config.max_wait_time}s")
                raise
            except WebDriverException as e:
                self.logger.error(f"WebDriver error: {e}")
                raise
    
    def _wait_for_content_load(self, driver: webdriver.Chrome) -> None:
        """Advanced waiting strategy for dynamic content"""
        wait_strategies = [
            # Wait for common ChatGPT canvas elements
            (By.CSS_SELECTOR, "[data-testid*='canvas']"),
            (By.CSS_SELECTOR, ".canvas-content"),
            (By.CSS_SELECTOR, "[role='main']"),
            (By.CLASS_NAME, "prose"),
            (By.TAG_NAME, "article"),
        ]
        
        for selector_type, selector in wait_strategies:
            try:
                WebDriverWait(driver, 5).until(
                    EC.presence_of_element_located((selector_type, selector))
                )
                self.logger.debug(f"Found content with selector: {selector}")
                break
            except TimeoutException:
                continue
        
        # Additional wait for JavaScript rendering
        time.sleep(3)
        
        # Check if page is still loading
        try:
            WebDriverWait(driver, 10).until(
                lambda d: d.execute_script("return document.readyState") == "complete"
            )
        except TimeoutException:
            self.logger.warning("Page may not be fully loaded")
    
    def _extract_content_multiple_strategies(self, soup: BeautifulSoup) -> Optional[Tag]:
        """Try multiple strategies to find the main content"""
        
        strategies = [
            self._extract_by_canvas_selectors,
            self._extract_by_semantic_tags,
            self._extract_by_content_density,
            self._extract_by_text_length,
        ]
        
        for strategy in strategies:
            try:
                content = strategy(soup)
                if content and self._validate_content(content):
                    self.logger.info(f"Content found using {strategy.__name__}")
                    return content
            except Exception as e:
                self.logger.debug(f"Strategy {strategy.__name__} failed: {e}")
        
        return None
    
    def _extract_by_canvas_selectors(self, soup: BeautifulSoup) -> Optional[Tag]:
        """Extract using ChatGPT-specific selectors"""
        selectors = [
            "[data-testid*='canvas']",
            ".canvas-content",
            "[role='main']",
            ".prose",
            "article",
            ".markdown",
        ]
        
        for selector in selectors:
            elements = soup.select(selector)
            for element in elements:
                if self._has_meaningful_content(element):
                    return element
        return None
    
    def _extract_by_semantic_tags(self, soup: BeautifulSoup) -> Optional[Tag]:
        """Extract using semantic HTML tags"""
        semantic_tags = ["main", "article", "section"]
        
        for tag_name in semantic_tags:
            elements = soup.find_all(tag_name)
            for element in elements:
                if self._has_meaningful_content(element):
                    return element
        return None
    
    def _extract_by_content_density(self, soup: BeautifulSoup) -> Optional[Tag]:
        """Find element with highest content density"""
        candidates = soup.find_all("div")
        best_candidate = None
        best_score = 0
        
        for div in candidates:
            score = self._calculate_content_score(div)
            if score > best_score:
                best_score = score
                best_candidate = div
        
        return best_candidate if best_score > 10 else None
    
    def _extract_by_text_length(self, soup: BeautifulSoup) -> Optional[Tag]:
        """Fallback: find element with most text content"""
        candidates = soup.find_all(["div", "section", "article", "main"])
        
        best_candidate = None
        max_length = 0
        
        for element in candidates:
            text_length = len(element.get_text(strip=True))
            if text_length > max_length and text_length > 100:
                max_length = text_length
                best_candidate = element
        
        return best_candidate
    
    def _calculate_content_score(self, element: Tag) -> int:
        """Calculate content richness score"""
        score = 0
        
        # Text content
        text = element.get_text(strip=True)
        score += len(text) // 10
        
        # Structural elements
        score += len(element.find_all(['h1', 'h2', 'h3'])) * 10
        score += len(element.find_all(['p'])) * 5
        score += len(element.find_all(['ul', 'ol'])) * 8
        score += len(element.find_all(['table'])) * 15
        score += len(element.find_all(['img'])) * 10
        
        # Penalty for too many generic divs (likely navigation/UI)
        score -= len(element.find_all(['div'])) // 5
        
        return max(0, score)
    
    def _has_meaningful_content(self, element: Tag) -> bool:
        """Check if element has meaningful content"""
        if not element:
            return False
        
        text = element.get_text(strip=True)
        if len(text) < 50:  # Too short to be meaningful
            return False
        
        # Check for structural elements
        structural_elements = element.find_all(['h1', 'h2', 'h3', 'p', 'ul', 'ol', 'table'])
        if len(structural_elements) < 2:
            return False
        
        return True
    
    def _validate_content(self, content: Tag) -> bool:
        """Validate extracted content quality"""
        if not content:
            return False
        
        text = content.get_text(strip=True)
        
        # Minimum content requirements
        if len(text) < 100:
            return False
        
        # Check for actual content structure
        structural_count = len(content.find_all(['h1', 'h2', 'h3', 'p', 'ul', 'ol', 'table', 'pre', 'code']))
        if structural_count < 3:
            return False
        
        return True


class TableProcessor:
    """Enhanced table processing with error handling"""
    
    def __init__(self, logger: logging.Logger):
        self.logger = logger
    
    def clean_and_validate_tables(self, soup: BeautifulSoup, content_div: Tag) -> Tag:
        """Clean tables with comprehensive error handling"""
        try:
            tables = content_div.find_all("table")
            self.logger.info(f"Processing {len(tables)} tables")
            
            for i, table in enumerate(tables):
                try:
                    self._clean_single_table(soup, table, i)
                except Exception as e:
                    self.logger.warning(f"Failed to clean table {i}: {e}")
                    # Remove problematic table rather than crash
                    table.decompose()
            
            return content_div
        except Exception as e:
            self.logger.error(f"Table processing failed: {e}")
            return content_div
    
    def _clean_single_table(self, soup: BeautifulSoup, table: Tag, table_index: int) -> None:
        """Clean a single table with validation"""
        rows = table.find_all("tr")
        if not rows:
            self.logger.warning(f"Table {table_index} has no rows")
            return
        
        # Calculate maximum columns safely
        max_cols = 0
        valid_rows = []
        
        for row in rows:
            cells = row.find_all(["td", "th"])
            if cells:  # Only process rows with cells
                max_cols = max(max_cols, len(cells))
                valid_rows.append((row, cells))
        
        if max_cols == 0:
            self.logger.warning(f"Table {table_index} has no valid cells")
            return
        
        # Normalize all rows to have the same number of columns
        for row, cells in valid_rows:
            while len(cells) < max_cols:
                new_cell = soup.new_tag("td")
                new_cell.string = ""
                row.append(new_cell)
                cells.append(new_cell)


class PowerPointGenerator:
    """Enhanced PowerPoint generation with advanced features"""
    
    def __init__(self, config: Config, logger: logging.Logger):
        self.speaker_notes_txt = [] 
        self.notes_seen = set()
        self.image_descriptions = []  # ADD THIS LINE
        self.config = config
        self.logger = logger
        self.slide_count = 0
    
    def create_enhanced_presentation(self, content_div: Tag, output_path: Path, title: str = None) -> bool:
        """Create PowerPoint with enhanced features and error handling"""
        try:
            prs = Presentation()
            
            self._set_default_fonts(prs)

            # Add title slide
            if title:
                self._add_title_slide(prs, title)
            
            # Process content elements
            self._process_content_elements(prs, content_div)
            
            # Ensure we have at least one slide
            if len(prs.slides) == 0:
                self._add_fallback_slide(prs, "No Content Found", "The canvas appears to be empty or could not be processed.")
            
            # Save with error handling
            self._save_presentation(prs, output_path)

            if self.speaker_notes_txt:
                self._save_speaker_notes_textfile(output_path, self.speaker_notes_txt)

            self.logger.info(f"✅ PowerPoint created with {len(prs.slides)} slides: {output_path}")
            
            return True
            
        except Exception as e:
            self.logger.error(f"PowerPoint generation failed: {e}")
            return False
    
    def _add_title_slide(self, prs: Presentation, title: str) -> None:
        """Add professional title slide"""
        try:
            title_slide_layout = prs.slide_layouts[0]  # Title Slide layout
            slide = prs.slides.add_slide(title_slide_layout)
            
            slide.shapes.title.text = title
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = f" "
            
            self.slide_count += 1
        except Exception as e:
            self.logger.warning(f"Failed to add title slide: {e}")
    
    def _process_content_elements(self, prs: Presentation, content_div: Tag) -> None:
        """Process all content elements with enhanced handling"""
        current_slide = None
        content_box = None
        
        # Get all relevant elements in document order
        elements = content_div.find_all([
            "h1", "h2", "h3", "h4", "h5", "h6",
            "p", "ul", "ol", "table", "pre", "code",
            "blockquote", "img", "span","div"
        ], recursive=True)
        
        code_buffer = []
        processed_elements = set()  # Track processed elements to avoid duplicates
    
        for element in elements:
            # Skip if already processed
            element_id = id(element)
            if element_id in processed_elements:
                continue
                
            element_text = element.get_text(strip=True)
            
            # Handle speaker notes
            if "Speaker notes:" in element_text:
                content_part, notes_part = re.split(r'Speaker notes\s*:\s*', element_text, flags=re.IGNORECASE, maxsplit=1)
    
                # If we have a current slide, add notes to it
                if current_slide is not None:
                    slide = current_slide
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = notes_part.strip()
                    notes_key = (self.slide_count, notes_part.strip())
                    if notes_key not in self.notes_seen:
                        self.notes_seen.add(notes_key)
                        self.speaker_notes_txt.append(notes_key)
                        notes_slide = slide.notes_slide
                        notes_slide.notes_text_frame.text = notes_part.strip()
                # Mark as processed and continue with content part if it exists
                processed_elements.add(element_id)
                if content_part.strip():
                    # Create a new element with just the content part for further processing
                    # But for now, we'll skip this element entirely after extracting notes
                    continue
                else:
                    continue
    
            # Handle image descriptions
            if "image:" in element_text.lower():
                content_part, image_part = re.split(r'image\s*:\s*', element_text, flags=re.IGNORECASE, maxsplit=1)

                # Store image description with current slide number if unique
                if image_part.strip():
                    desc = image_part.strip()
                    slide_no = self.slide_count + 1
                    if not any(d['slide_number'] == slide_no and d['description'] == desc for d in self.image_descriptions):
                        self.image_descriptions.append({
                            'slide_number': slide_no,  # Next slide number
                            'description': desc
                        })
                
                # Mark as processed
                processed_elements.add(element_id)
                
                # Continue processing the content part if it exists
                if content_part.strip():
                    # For now, skip further processing to avoid duplication
                    # You could create a new element here if needed
                    continue
                else:
                    continue
    
            # Handle consecutive cm-line blocks as one code block
            if element.name == "div" and "cm-line" in element.get("class", []):
                print(f"[DEBUG] Detected cm-line: {element.get_text(strip=True)}")
                code_text = element.get_text(strip=True)
                if code_text:
                    code_buffer.append(code_text)
                processed_elements.add(element_id)
                continue
            
            # If the current element is NOT a cm-line AND we have code buffered, flush it
            if code_buffer:
                current_slide, content_box = self._ensure_slide(prs, current_slide, "Content")
                self._add_code_content(content_box, "\n".join(code_buffer))
                code_buffer = []
            
            # Skip elements that are part of other elements to avoid duplication
            if element.name in ["p", "span"] and element.find_parents(["ul", "ol", "li"]):
                processed_elements.add(element_id)
                continue
            
            if element.name == "p" and element.find_parent("li"):
                processed_elements.add(element_id)
                continue
            
            if element.name == "span" and (element.find_parent("p") or element.find_parent("li")):
                processed_elements.add(element_id)
                continue
    
            # Skip <code> if it's inside a <pre>
            if element.name == "code" and element.find_parent("pre"):
                processed_elements.add(element_id)
                continue
    
            try:
                element_type = element.name
                
                # Handle headings - create new slides
                if element_type in ["h1", "h2", "h3", "h4", "h5", "h6"]:
                    current_slide, content_box = self._add_content_slide(prs, element.get_text(strip=True))
                
                # Handle lists
                elif element_type in ["ol", "ul"] and not element.find_parent(["ul", "ol"]):
                    current_slide, content_box = self._ensure_slide(prs, current_slide, "List")
                    self._add_list_content(content_box, element)
                    processed_elements.add(element_id)
                    continue
                
                # Handle tables
                elif element_type == "table":
                   current_slide, content_box = self._ensure_slide(prs, current_slide, "Content")
                   self._add_table_to_slide(current_slide, element)
                   processed_elements.add(element_id)
                   continue
                
                # Handle code blocks
                elif element_type in ["pre", "code"]:
                    current_slide, content_box = self._ensure_slide(prs, current_slide, "Code")
                    self._add_code_content(content_box, element)
                
                # Mark as processed
                processed_elements.add(element_id)
                    
            except Exception as e:
                self.logger.warning(f"Failed to process element {element_type}: {e}")
                processed_elements.add(element_id)
                continue
            
            # Handle any remaining code buffer
            if code_buffer:
                current_slide, content_box = self._ensure_slide(prs, current_slide, "Content")
                self._add_code_content(content_box, "\n".join(code_buffer))

    def _add_content_slide(self, prs: Presentation, title: str) -> Tuple[Any, Any]:
        """Add a new content slide"""
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = re.sub(r'^slide\s*\d+\s*:\s*', '', title, flags=re.IGNORECASE)

        # Set title with length check
        slide.shapes.title.text = title[:100] + "..." if len(title) > 100 else title
        
        # Get content placeholder
        content_box = slide.placeholders[1]
        content_box.text_frame.clear()
        
        self.slide_count += 1
        return slide, content_box
    
    def _ensure_slide(self, prs: Presentation, current_slide: Any, default_title: str) -> Tuple[Any, Any]:
        """Ensure we have a slide to work with"""
        if current_slide is None:
            return self._add_content_slide(prs, default_title)
        
        # Return existing slide and its content box
        content_box = current_slide.placeholders[1]
        return current_slide, content_box
    
    def _add_paragraph_content(self, content_box: Any, element: Tag) -> None:
        """Add paragraph with smart formatting"""
        print("Para content being called")
        text = element.get_text(strip=True)
        if not text or len(text) > self.config.max_slide_content_length:
            return
        if "speaker notes:" in text.lower():
           content_part, notes_part = re.split(r'speaker notes\s*:\s*', text, flags=re.IGNORECASE, maxsplit=1)
           text = content_part.strip()
       
           # Add speaker notes to the slide
           notes_slide = content_box.part.slide.notes_slide
           notes_slide.notes_text_frame.text = notes_part.strip()

        # Check if this looks like a bullet point
        bullet_patterns = [r'^\s*[-•·]\s+', r'^\s*\d+\.\s+', r'^\s*[a-zA-Z]\.\s+']
        is_bullet = any(re.match(pattern, text) for pattern in bullet_patterns)
        
        para = content_box.text_frame.add_paragraph()
        para.text = text
        para.level = 1 if is_bullet else 0
        
        
        self._set_font_safely(para, text, 'default')
    
    def _add_list_content(self, content_box: Any, list_element: Tag) -> None:
        """Add list with proper hierarchy"""
        self._process_list_recursive(content_box, list_element, 0)
    
    def _process_list_recursive(self, content_box: Any, list_element: Tag, level: int) -> None:
        """Process lists recursively with proper nesting"""
        max_level = 4  # PowerPoint limitation
        print("rec list being called")
        for li in list_element.find_all("li", recursive=False):
            try:
                # Get text content, excluding nested lists
                text_parts = []
                for item in li.children:
                    if isinstance(item, NavigableString):
                        text_parts.append(str(item).strip())
                    elif isinstance(item, Tag) and item.name not in ["ul", "ol"]:
                        text_parts.append(item.get_text(" ", strip=True))
                
                text = " ".join(text_parts).strip()
                if text:
                    para = content_box.text_frame.add_paragraph()
                    para.text = text
                    para.level = min(level, max_level)
                    
                    self._set_font_safely(para, text, 'default')
                
                # Process nested lists
                nested_lists = li.find_all(["ul", "ol"], recursive=False)
                for nested in nested_lists:
                    self._process_list_recursive(content_box, nested, level + 1)
                    
            except Exception as e:
                self.logger.debug(f"List item processing failed: {e}")
    
    def _add_table_to_slide(self, slide: Any, table_element: Tag) -> None:
        """Insert table into an existing slide"""
        try:
            rows = table_element.find_all("tr")
            if not rows:
                return
            
            # Calculate table dimensions
            max_cols = max(len(row.find_all(["td", "th"])) for row in rows if row.find_all(["td", "th"]))
            num_rows = len(rows)
            
            if max_cols == 0 or num_rows == 0:
                return
            
            # Table dimensions and positioning
            left = Inches(0.5)
            top = Inches(3.5)  # Position below title or existing content
            width = Inches(9)
            height = Inches(min(5.5, 0.5 + 0.4 * num_rows))
            
            # Create table
            table_shape = slide.shapes.add_table(num_rows, max_cols, left, top, width, height)
            table = table_shape.table
            
            for i, row in enumerate(rows):
                cells = row.find_all(["td", "th"])
                is_header = any(cell.name == "th" for cell in cells)
                
                for j in range(max_cols):
                    cell_text = ""
                    if j < len(cells):
                        cell_text = cells[j].get_text(strip=True)
                    
                    cell = table.cell(i, j)
                    cell.text = cell_text[:200]
                    
                    if is_header and i == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = RGBColor(255, 255, 255)
                                run.font.bold = True
        except Exception as e:
            self.logger.warning(f"Table insertion failed: {e}")

    
    def _add_code_content(self, content_box: Any, code_text: str) -> None:
        """Add code with monospace formatting and no bullets"""
        if not code_text.strip():
            return
    
        para = content_box.text_frame.add_paragraph()
        para.text = code_text[:self.config.max_slide_content_length]
        para.level = 0  # Make sure it's top-level
    
        # 🔧 Safely remove bullets
        pPr = para._element.get_or_add_pPr()
        for bullet_tag in ['a:buAutoNum', 'a:buChar', 'a:buNone']:
            tag = pPr.find(qn(bullet_tag))
            if tag is not None:
                pPr.remove(tag)
    
        # ⛑️ Add bullet=None explicitly (if using a style that enforces bullets)
        buNone = OxmlElement('a:buNone')
        pPr.append(buNone)
    
        self._set_font_safely(para, code_text, 'code')

    
    def _add_quote_content(self, content_box: Any, element: Tag) -> None:
        """Add blockquote with special formatting"""
        quote_text = element.get_text(strip=True)
        if not quote_text:
            return
        
        para = content_box.text_frame.add_paragraph()
        para.text = f'"{quote_text}"'
        para.level = 0
        
        self._set_font_safely(para, quote_text, 'default')
        try:
            para.font.italic = True
        except Exception:
            pass
    
    # '''def _add_math_content(self, content_box: Any, element: Tag) -> None:
    #     """Add mathematical expressions"""
    #     math_text = element.get_text(strip=True)
    #     if not math_text:
    #         return
        
    #     para = content_box.text_frame.add_paragraph()
    #     para.text = f"Formula: {math_text}"
        
    #     try:
    #         para.font.name = self.config.font_fallbacks['math']
    #         para.font.size = Pt(26)
    #     except Exception:
    #         pass
    
    # def _is_math_element(self, element: Tag) -> bool:
    #     """Check if element contains mathematical expressions"""
    #     class_names = element.get("class", [])
    #     return any("katex" in str(cls).lower() or "math" in str(cls).lower() for cls in class_names)'''
    
    # def _get_appropriate_font(self, text: str) -> str:
    #     """Get appropriate font based on text content"""
    #     # Check for non-ASCII characters (might need special font handling)
    #     if any(ord(c) > 127 for c in text):
    #         return self.config.font_fallbacks['fallback']
        
    #     # Check for code-like content
    #     if re.search(r'[{}();=<>]', text) and len(text.split()) < 10:
    #         return self.config.font_fallbacks['code']
        
    #     return self.config.font_fallbacks['default']
    
    def _add_fallback_slide(self, prs: Presentation, title: str, content: str) -> None:
        """Add fallback slide when no content is found"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
        self.slide_count += 1

    def _set_font_safely(self, paragraph, text_content, font_type='default'):
        """Safely set font with proper error handling and logging"""
        try:
            if font_type == 'code':
                font_name = self.config.font_fallbacks['code']
                font_size = Pt(20)
            elif font_type == 'heading':
                font_name = self.config.font_fallbacks['calibri']
                font_size = Pt(28)
            else:
                font_name = self._get_appropriate_font(text_content)
                font_size = Pt(24)
            
            paragraph.font.name = font_name
            paragraph.font.size = font_size
            
            self.logger.debug(f"Font set successfully: {font_name}, {font_size}")
            
        except Exception as e:
            self.logger.warning(f"Font setting failed for '{text_content[:50]}...': {e}")
            # Try fallback
            try:
                paragraph.font.name = 'Arial'
                paragraph.font.size = Pt(24)
            except Exception as fallback_error:
                self.logger.error(f"Even fallback font failed: {fallback_error}")
    
    def _set_default_fonts(self, prs: Presentation):
        """Set consistent default fonts across all slide layouts"""
        try:
            for layout in prs.slide_layouts:
                for placeholder in layout.placeholders:
                    if hasattr(placeholder, 'text_frame'):
                        for paragraph in placeholder.text_frame.paragraphs:
                            try:
                                paragraph.font.name = self.config.font_fallbacks['default']
                                paragraph.font.size = Pt(24)
                            except:
                                continue
        except Exception as e:
            self.logger.debug(f"Default font setting failed: {e}")

    def _save_presentation(self, prs: Presentation, output_path: Path) -> None:
        """Save presentation with comprehensive error handling"""
        try:
            # Ensure directory exists
            output_path.parent.mkdir(parents=True, exist_ok=True)
            
            # Check available disk space (basic check)
            if hasattr(os, 'statvfs'):  # Unix systems
                stat = os.statvfs(output_path.parent)
                available_space = stat.f_frsize * stat.f_available
                if available_space < 10 * 1024 * 1024:  # Less than 10MB
                    raise IOError("Insufficient disk space")
            
            # Save presentation
            prs.save(str(output_path))
            
        except PermissionError:
            raise IOError(f"Permission denied: Cannot write to {output_path}")
        except OSError as e:
            raise IOError(f"File system error: {e}")
    
    def _create_image_document(self, ppt_path: Path) -> None:
        """Create a Word document with image descriptions"""
        try:
            doc = Document()
            doc.add_heading('Image Descriptions', 0)
            # doc.add_paragraph(f'Generated from PowerPoint: {ppt_path.name}')
            # doc.add_paragraph(f'Created on: {time.strftime("%Y-%m-%d %H:%M:%S")}')
            doc.add_paragraph('')  # Empty line
            
            seen = set()
            unique_descriptions = []
            for img_info in self.image_descriptions:
                key = (img_info.get('slide_number'), img_info.get('description'))
                if key not in seen:
                    seen.add(key)
                    unique_descriptions.append(img_info)

            for img_info in unique_descriptions:
                # Add slide number as heading
                doc.add_heading(f'Slide {img_info["slide_number"]}', level=1)
                # Add image description
                doc.add_paragraph(img_info['description'])
                doc.add_paragraph('')  # Empty line between entries
            
            # Save document with same name as PowerPoint but .docx extension
            doc_path = ppt_path.with_suffix('.docx')
            doc_path = SafeFilename.ensure_unique(doc_path)
            doc.save(str(doc_path))
            
            self.logger.info(f"📄 Image descriptions document created: {doc_path}")
            
        except Exception as e:
            self.logger.error(f"Failed to create image descriptions document: {e}")
    def _save_speaker_notes_textfile(self, ppt_path: Path, speaker_notes_list: List[Tuple[int, str]]) -> None:
        """Save speaker notes to a text file with spacing between slides"""
        try:
            textfile_path = ppt_path.with_name(ppt_path.stem + "_speaker_notes.txt")
            with open(textfile_path, "w", encoding="utf-8") as f:
                for slide_number, notes in speaker_notes_list:
                    f.write(f"Slide {slide_number}:\n{notes.strip()}\n\n")
            self.logger.info(f"📝 Speaker notes text file saved: {textfile_path}")
        except Exception as e:
            self.logger.error(f"Failed to save speaker notes text file: {e}")
    
class CanvasConverter:
    """Main converter class orchestrating the entire process"""
    
    def __init__(self, config: Config = None, log_level: str = 'INFO'):
        self.config = config or Config()
        self.logger = setup_logging(log_level)
        self.extractor = ContentExtractor(self.config, self.logger)
        self.table_processor = TableProcessor(self.logger)
        self.ppt_generator = PowerPointGenerator(self.config, self.logger)
    
    def convert(self, url: str, output_dir: str = None, filename: str = None) -> Optional[Path]:
        """Main conversion method with comprehensive error handling"""
        try:
            self.logger.info("🚀 Starting ChatGPT Canvas to PowerPoint conversion")
            
            # Step 1: Extract content
            self.logger.info("🔍 Extracting content from canvas...")
            soup, content = self.extractor.scrape_with_retry(url)
            
            if not content:
                self.logger.error("❌ Failed to extract content from canvas")
                return None
            
            # Step 2: Process tables
            self.logger.info("🧹 Processing tables...")
            processed_content = self.table_processor.clean_and_validate_tables(soup, content)
            
            # Step 3: Generate filename
            output_path = self._generate_output_path(soup, output_dir, filename)
            
            # Step 4: Create PowerPoint
            self.logger.info("🖼️  Generating PowerPoint presentation...")
            title = self._extract_title(soup)
            
            success = self.ppt_generator.create_enhanced_presentation(
                processed_content, output_path, title
            )
            
            if success:
                self.logger.info(f"✅ Conversion completed successfully: {output_path}")
                return output_path
            else:
                self.logger.error("❌ PowerPoint generation failed")
                return None
                
        except Exception as e:
            self.logger.error(f"❌ Conversion failed: {e}")
            return None
    
    def _generate_output_path(self, soup: BeautifulSoup, output_dir: str = None, filename: str = None) -> Path:
        """Generate safe output path"""
        if filename:
            base_name = SafeFilename.sanitize(filename, self.config.max_filename_length)
        else:
            # Extract title from page
            title = self._extract_title(soup)
            base_name = SafeFilename.sanitize(title, self.config.max_filename_length)
        
        if not base_name.endswith('.pptx'):
            base_name += '.pptx'
        
        if output_dir:
            output_path = Path(output_dir) / base_name
        else:
            output_path = Path.cwd() / base_name
        
        return SafeFilename.ensure_unique(output_path)
    
    def _extract_title(self, soup: BeautifulSoup) -> str:
        """Extract title from page with multiple strategies"""
        # Try different title extraction methods
        title_candidates = [
            soup.find("h1"),
            soup.find("title"),
            soup.find(attrs={"data-testid": "canvas-title"}),
            soup.find(class_=re.compile("title", re.I)),
        ]
        
        for candidate in title_candidates:
            if candidate and candidate.get_text(strip=True):
                return candidate.get_text(strip=True)
        
        return "ChatGPT_Canvas_Export"
    
    def batch_convert(self, urls: List[str], output_dir: str = None) -> Dict[str, Optional[Path]]:
        """Convert multiple canvas URLs"""
        results = {}
        
        for i, url in enumerate(urls, 1):
            self.logger.info(f"Processing {i}/{len(urls)}: {url}")
            try:
                result = self.convert(url, output_dir)
                results[url] = result
            except Exception as e:
                self.logger.error(f"Failed to process {url}: {e}")
                results[url] = None
        
        return results


class CLIInterface:
    """Command-line interface for the converter"""
    
    def __init__(self):
        self.config = Config()
        self.converter = CanvasConverter(self.config)
    
    def run_interactive(self):
        """Run interactive CLI"""
        print("🎯 Enhanced ChatGPT Canvas to PowerPoint Converter")
        print("=" * 60)
        
        while True:
            try:
                # Get URL
                url = input("\n🔗 Enter ChatGPT canvas share URL (or 'quit' to exit): ").strip()
                
                if url.lower() in ['quit', 'exit', 'q']:
                    print("👋 Goodbye!")
                    break
                
                if not url:
                    print("❌ Please enter a valid URL")
                    continue
                
                # Get optional parameters
                output_dir = input("📁 Output directory (press Enter for current): ").strip() or None
                filename = input("📄 Custom filename (press Enter for auto): ").strip() or None
                
                # Convert
                result = self.converter.convert(url, output_dir, filename)
                
                if result:
                    print(f"\n✅ Success! PowerPoint saved to: {result}")
                    
                    # Ask if user wants to open the file
                    if self._ask_yes_no("🚀 Open the PowerPoint file now?"):
                        self._open_file(result)
                else:
                    print("\n❌ Conversion failed. Check the logs above for details.")
                
            except KeyboardInterrupt:
                print("\n\n👋 Interrupted by user. Goodbye!")
                break
            except Exception as e:
                print(f"\n❌ Unexpected error: {e}")
    
    def run_batch(self, urls_file: str, output_dir: str = None):
        """Run batch conversion from file"""
        try:
            with open(urls_file, 'r') as f:
                urls = [line.strip() for line in f if line.strip() and not line.startswith('#')]
            
            if not urls:
                print("❌ No valid URLs found in file")
                return
            
            print(f"🔄 Processing {len(urls)} URLs...")
            results = self.converter.batch_convert(urls, output_dir)
            
            # Summary
            successful = sum(1 for result in results.values() if result)
            print(f"\n📊 Batch conversion complete:")
            print(f"   ✅ Successful: {successful}")
            print(f"   ❌ Failed: {len(results) - successful}")
            print(f"   📁 Output directory: {output_dir or 'current directory'}")
            
            ppt_paths = [p for p in results.values() if p is not None]
            if len(ppt_paths) >= 2:
                print("🔗 Merging presentations...")
                final_merged = merge_presentations(ppt_paths[0], ppt_paths[1:])
                print(f"📦 Combined PPT saved at: {final_merged}")

        except FileNotFoundError:
            print(f"❌ File not found: {urls_file}")
        except Exception as e:
            print(f"❌ Batch processing failed: {e}")
    
    def _ask_yes_no(self, question: str) -> bool:
        """Ask yes/no question"""
        while True:
            answer = input(f"{question} (y/n): ").strip().lower()
            if answer in ['y', 'yes']:
                return True
            elif answer in ['n', 'no']:
                return False
            else:
                print("Please enter 'y' or 'n'")
    
    def _open_file(self, filepath: Path):
        """Open file with system default application"""
        try:
            if sys.platform.startswith('darwin'):  # macOS
                os.system(f'open "{filepath}"')
            elif sys.platform.startswith('win'):  # Windows
                os.startfile(str(filepath))
            else:  # Linux and others
                os.system(f'xdg-open "{filepath}"')
        except Exception as e:
            print(f"⚠️  Could not open file automatically: {e}")
            print(f"📁 File location: {filepath}")


def main():
    """Main entry point"""
    import argparse
    
    parser = argparse.ArgumentParser(
        description="Enhanced ChatGPT Canvas to PowerPoint Converter",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python canvas_converter.py                          # Interactive mode
  python canvas_converter.py -u "https://..."        # Single conversion
  python canvas_converter.py -b urls.txt -o output/  # Batch conversion
  python canvas_converter.py -u "https://..." -v     # Verbose logging
        """
    )
    
    parser.add_argument("-u", "--url", help="ChatGPT canvas URL to convert")
    parser.add_argument("-o", "--output", help="Output directory")
    parser.add_argument("-f", "--filename", help="Custom output filename")
    parser.add_argument("-b", "--batch", help="File containing URLs for batch processing")
    parser.add_argument("-v", "--verbose", action="store_true", help="Verbose logging")
    parser.add_argument("--log-level", choices=['DEBUG', 'INFO', 'WARNING', 'ERROR'], 
                       default='INFO', help="Logging level")
    parser.add_argument("--merge", action="store_true", help="Merge all generated PPTs into one")

    
    args = parser.parse_args()
    
    # Set log level
    log_level = 'DEBUG' if args.verbose else args.log_level
    
    # Create CLI interface
    cli = CLIInterface()
    cli.converter = CanvasConverter(log_level=log_level)
    
    try:
        if args.batch:
            # Batch mode
            cli.run_batch(args.batch, args.output)
        elif args.url:
            # Single URL mode
            result = cli.converter.convert(args.url, args.output, args.filename)
            if result:
                print(f"✅ Success! PowerPoint saved to: {result}")
            else:
                print("❌ Conversion failed")
                sys.exit(1)
        else:
            # Interactive mode
            cli.run_interactive()
            
    except KeyboardInterrupt:
        print("\n👋 Interrupted by user")
    except Exception as e:
        print(f"❌ Fatal error: {e}")
        sys.exit(1)

def merge_presentations(primary_path: Path, additional_paths: List[Path], output_path: Path = None) -> Path:
    """
    Merges additional PowerPoint files into a primary one.

    :param primary_path: Path to the base presentation
    :param additional_paths: List of presentations to merge in order
    :param output_path: Path to save the merged presentation
    :return: Path to the final merged presentation
    """
    merged = Presentation(primary_path)

    for ppt_path in additional_paths:
        to_merge = Presentation(ppt_path)

        for slide in to_merge.slides:
            # Copy slide
            slide_layout = merged.slide_layouts[1]  # Use content layout
            new_slide = merged.slides.add_slide(slide_layout)

            for shape in slide.shapes:
                try:
                    el = shape.element
                    new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
                except Exception as e:
                    print(f"⚠️ Error copying shape: {e}")

            # Copy notes
            try:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text:
                    notes = new_slide.notes_slide.notes_text_frame
                    notes.text = notes_text
            except Exception:
                pass

    final_path = output_path or primary_path.with_stem(primary_path.stem + "_merged")
    final_path = SafeFilename.ensure_unique(final_path.with_suffix('.pptx'))
    merged.save(final_path)
    return final_path

if __name__ == "__main__":
    main()