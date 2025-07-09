#!/usr/bin/env python3
"""
PowerPoint Image Inserter Script with Document Mapping
Inserts images into PowerPoint slides based on mappings from a configuration document
"""

import os
import sys
import json
import csv
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Union
from pptx import Presentation
from pptx.util import Inches, Cm
from pptx.enum.shapes import MSO_SHAPE_TYPE

class PPTImageInserter:
    def __init__(self):
        self.image_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp'}
        self.image_files = []
        
    def get_image_files(self, directory: str = ".") -> List[Path]:
        """Get list of image files in directory"""
        image_files = []
        directory_path = Path(directory)
        
        if not directory_path.exists():
            raise FileNotFoundError(f"Image directory '{directory}' not found")
            
        for file_path in directory_path.glob("*"):
            if file_path.suffix.lower() in self.image_extensions:
                image_files.append(file_path)
        
        return sorted(image_files)
    
    def parse_mapping_document(self, mapping_file: str) -> List[Dict]:
        """Parse mapping document (supports JSON, CSV, or TXT formats)"""
        mapping_path = Path(mapping_file)
        
        if not mapping_path.exists():
            raise FileNotFoundError(f"Mapping file '{mapping_file}' not found")
        
        file_extension = mapping_path.suffix.lower()
        
        if file_extension == '.json':
            return self._parse_json_mapping(mapping_path)
        elif file_extension == '.csv':
            return self._parse_csv_mapping(mapping_path)
        elif file_extension in ['.txt', '.map']:
            return self._parse_txt_mapping(mapping_path)
        else:
            raise ValueError(f"Unsupported mapping file format: {file_extension}")
    
    def _parse_json_mapping(self, file_path: Path) -> List[Dict]:
        """Parse JSON mapping file
        Expected format:
        [
            {
                "image_number": 1,
                "slide_number": 1,
                "position": "center",
                "left": 1.0,
                "top": 1.0,
                "width": 4.0,
                "height": 3.0
            }
        ]
        """
        with open(file_path, 'r') as f:
            data = json.load(f)
        
        mappings = []
        for item in data:
            mapping = {
                'image_number': item.get('image_number'),
                'slide_number': item.get('slide_number'),
                'position': item.get('position', 'center'),
                'left': item.get('left'),
                'top': item.get('top'),
                'width': item.get('width'),
                'height': item.get('height')
            }
            mappings.append(mapping)
        
        return mappings
    
    def _parse_csv_mapping(self, file_path: Path) -> List[Dict]:
        """Parse CSV mapping file
        Expected columns: image_number, slide_number, position, left, top, width, height
        """
        mappings = []
        
        with open(file_path, 'r', newline='') as f:
            reader = csv.DictReader(f)
            for row in reader:
                mapping = {
                    'image_number': int(row['image_number']) if row.get('image_number') else None,
                    'slide_number': int(row['slide_number']) if row.get('slide_number') else None,
                    'position': row.get('position', 'center'),
                    'left': float(row['left']) if row.get('left') and row['left'].strip() else None,
                    'top': float(row['top']) if row.get('top') and row['top'].strip() else None,
                    'width': float(row['width']) if row.get('width') and row['width'].strip() else None,
                    'height': float(row['height']) if row.get('height') and row['height'].strip() else None
                }
                mappings.append(mapping)
        
        return mappings
    
    def _parse_txt_mapping(self, file_path: Path) -> List[Dict]:
        """Parse TXT mapping file
        Expected format (one mapping per line):
        image_number:slide_number:position:left:top:width:height
        or
        image_number slide_number position left top width height
        """
        mappings = []
        
        with open(file_path, 'r') as f:
            for line_num, line in enumerate(f, 1):
                line = line.strip()
                if not line or line.startswith('#'):  # Skip empty lines and comments
                    continue
                
                # Try colon-separated first, then space-separated
                if ':' in line:
                    parts = line.split(':')
                else:
                    parts = line.split()
                
                if len(parts) < 2:
                    print(f"Warning: Skipping invalid line {line_num}: {line}")
                    continue
                
                try:
                    mapping = {
                        'image_number': int(parts[0]) if parts[0].strip() else None,
                        'slide_number': int(parts[1]) if parts[1].strip() else None,
                        'position': parts[2].strip() if len(parts) > 2 and parts[2].strip() else 'center',
                        'left': float(parts[3]) if len(parts) > 3 and parts[3].strip() else None,
                        'top': float(parts[4]) if len(parts) > 4 and parts[4].strip() else None,
                        'width': float(parts[5]) if len(parts) > 5 and parts[5].strip() else None,
                        'height': float(parts[6]) if len(parts) > 6 and parts[6].strip() else None
                    }
                    mappings.append(mapping)
                except (ValueError, IndexError) as e:
                    print(f"Warning: Error parsing line {line_num}: {line} - {e}")
                    continue
        
        return mappings
    
    def get_position_settings(self, mapping: Dict) -> Tuple[str, Optional[float], Optional[float], Optional[float], Optional[float]]:
        """Convert mapping to position settings"""
        position = mapping.get('position', 'center')
        left = Inches(mapping['left']) if mapping.get('left') is not None else None
        top = Inches(mapping['top']) if mapping.get('top') is not None else None
        width = Inches(mapping['width']) if mapping.get('width') is not None else None
        height = Inches(mapping['height']) if mapping.get('height') is not None else None
        
        # Handle preset positions
        if position == 'top-left':
            left, top = Inches(0.5), Inches(0.5)
        elif position == 'top-right':
            left, top = Inches(6), Inches(0.5)
        elif position == 'bottom-left':
            left, top = Inches(0.5), Inches(5)
        elif position == 'bottom-right':
            left, top = Inches(6), Inches(5)
        elif position == 'center' and (left is None or top is None):
            return 'center', None, None, width, height
        
        return 'custom', left, top, width, height
    
    def insert_image_to_slide(self, slide, image_path: Path, position_type: str, 
                            left=None, top=None, width=None, height=None) -> bool:
        """Insert image into a slide at specified position"""
        try:
            if position_type == "center":
                # Calculate center position
                slide_width = Inches(10)  # Standard slide width
                slide_height = Inches(7.5)  # Standard slide height
                
                # Add image first to get its natural dimensions
                temp_pic = slide.shapes.add_picture(str(image_path), Inches(0), Inches(0), 
                                                  width=width, height=height)
                
                # Calculate centered position
                img_width = temp_pic.width
                img_height = temp_pic.height
                
                # Remove temporary image
                sp = temp_pic._element
                sp.getparent().remove(sp)
                
                # Calculate center position
                left = (slide_width - img_width) / 2
                top = (slide_height - img_height) / 2
                
                # Add image at center
                picture = slide.shapes.add_picture(str(image_path), left, top, img_width, img_height)
            else:
                # Custom position
                if width and height:
                    picture = slide.shapes.add_picture(str(image_path), left, top, width, height)
                elif width:
                    picture = slide.shapes.add_picture(str(image_path), left, top, width=width)
                elif height:
                    picture = slide.shapes.add_picture(str(image_path), left, top, height=height)
                else:
                    picture = slide.shapes.add_picture(str(image_path), left, top)
            
            return True
            
        except Exception as e:
            print(f"✗ Error inserting image {image_path.name}: {e}")
            return False
    
    def process_mappings(self, ppt_file: str, image_dir: str, mapping_file: str, output_file: str = None):
        """Main processing function"""
        print("=== PowerPoint Image Inserter with Document Mapping ===\n")
        
        # Load or create presentation
        if os.path.exists(ppt_file):
            try:
                prs = Presentation(ppt_file)
                print(f"✓ Loaded presentation: {ppt_file}")
            except Exception as e:
                print(f"✗ Error loading presentation: {e}")
                return False
        else:
            print(f"File '{ppt_file}' not found. Creating new presentation...")
            prs = Presentation()
            # Add a blank slide
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            prs.slides.add_slide(blank_slide_layout)
            print("✓ Created new presentation with 1 blank slide.")
        
        print(f"Presentation has {len(prs.slides)} slides.\n")
        
        # Get image files
        try:
            self.image_files = self.get_image_files(image_dir)
            print(f"✓ Found {len(self.image_files)} image files in '{image_dir}'")
            for i, img_file in enumerate(self.image_files, 1):
                print(f"  {i}. {img_file.name}")
            print()
        except Exception as e:
            print(f"✗ Error loading images: {e}")
            return False
        
        # Parse mapping document
        try:
            mappings = self.parse_mapping_document(mapping_file)
            print(f"✓ Parsed {len(mappings)} mappings from '{mapping_file}'\n")
        except Exception as e:
            print(f"✗ Error parsing mapping file: {e}")
            return False
        
        # Process each mapping
        success_count = 0
        for i, mapping in enumerate(mappings, 1):
            print(f"Processing mapping {i}/{len(mappings)}:")
            print(f"  Image {mapping['image_number']} → Slide {mapping['slide_number']}")
            
            # Validate mapping
            if not mapping['image_number'] or not mapping['slide_number']:
                print("  ✗ Invalid mapping: missing image_number or slide_number")
                continue
            
            if mapping['image_number'] < 1 or mapping['image_number'] > len(self.image_files):
                print(f"  ✗ Invalid image number: {mapping['image_number']}")
                continue
            
            if mapping['slide_number'] < 1 or mapping['slide_number'] > len(prs.slides):
                print(f"  ✗ Invalid slide number: {mapping['slide_number']}")
                continue
            
            # Get image and slide
            image_file = self.image_files[mapping['image_number'] - 1]
            slide = prs.slides[mapping['slide_number'] - 1]
            
            # Get position settings
            position_type, left, top, width, height = self.get_position_settings(mapping)
            
            # Insert image
            if self.insert_image_to_slide(slide, image_file, position_type, left, top, width, height):
                print(f"  ✓ Inserted {image_file.name} into slide {mapping['slide_number']}")
                success_count += 1
            else:
                print(f"  ✗ Failed to insert {image_file.name}")
        
        # Save presentation
        try:
            if not output_file:
                output_file = ppt_file
            prs.save(output_file)
            print(f"\n✓ Successfully processed {success_count}/{len(mappings)} mappings")
            print(f"✓ Presentation saved as: {output_file}")
            return True
        except Exception as e:
            print(f"\n✗ Error saving presentation: {e}")
            return False

def create_sample_mapping_files():
    """Create sample mapping files for reference"""
    
    # Sample JSON mapping
    json_sample = [
        {"image_number": 1, "slide_number": 1, "position": "center"},
        {"image_number": 2, "slide_number": 2, "position": "top-left"},
        {"image_number": 3, "slide_number": 3, "position": "custom", "left": 2.0, "top": 1.5, "width": 5.0, "height": 4.0}
    ]
    
    with open('sample_mapping.json', 'w') as f:
        json.dump(json_sample, f, indent=2)
    
    # Sample CSV mapping
    csv_content = """image_number,slide_number,position,left,top,width,height
1,1,center,,,4.0,3.0
2,2,top-left,,,,
3,3,custom,2.0,1.5,5.0,4.0
4,4,bottom-right,,,,"""
    
    with open('sample_mapping.csv', 'w') as f:
        f.write(csv_content)
    
    # Sample TXT mapping
    txt_content = """# Image mapping file
# Format: image_number:slide_number:position:left:top:width:height
# Position can be: center, top-left, top-right, bottom-left, bottom-right, custom
1:1:center:::4.0:3.0
2:2:top-left::::
3:3:custom:2.0:1.5:5.0:4.0
4:4:bottom-right::::"""
    
    with open('sample_mapping.txt', 'w') as f:
        f.write(txt_content)
    
    print("✓ Created sample mapping files: sample_mapping.json, sample_mapping.csv, sample_mapping.txt")

def main():
    if len(sys.argv) > 1 and sys.argv[1] == '--create-samples':
        create_sample_mapping_files()
        return
    
    inserter = PPTImageInserter()
    
    try:
        # Get inputs
        ppt_file = input("Enter PowerPoint file path (or press Enter for 'presentation.pptx'): ").strip()
        if not ppt_file:
            ppt_file = "presentation.pptx"
        
        image_dir = input("Enter images directory (press Enter for current directory): ").strip()
        if not image_dir:
            image_dir = "."
        
        mapping_file = input("Enter mapping file path (JSON/CSV/TXT): ").strip()
        if not mapping_file:
            print("Mapping file is required!")
            return
        
        output_file = input(f"Save as (press Enter for '{ppt_file}'): ").strip()
        if not output_file:
            output_file = ppt_file
        
        # Process mappings
        inserter.process_mappings(ppt_file, image_dir, mapping_file, output_file)
        
    except KeyboardInterrupt:
        print("\n\nOperation cancelled by user.")
    except Exception as e:
        print(f"\nUnexpected error: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()